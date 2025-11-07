import os
import streamlit as st
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

# ✅ LangChain imports for v1.0.4 + community 0.4.1
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.prompts import PromptTemplate
from langchain.chains.combine_documents.stuff import create_stuff_documents_chain
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.combine_documents.stuff import StuffDocumentsChain
from langchain.chains.llm import LLMChain

# --- Load environment variables ---
load_dotenv()
if "GOOGLE_API_KEY" not in os.environ or not os.environ["GOOGLE_API_KEY"]:
    st.error("❌ GOOGLE_API_KEY not found in .env file.")
    st.stop()


# --- Extract text from PDF, DOCX, PPTX ---
def get_text_from_files(files):
    text = ""
    for file in files:
        ext = os.path.splitext(file.name)[1].lower()

        if ext == ".pdf":
            reader = PdfReader(file)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"

        elif ext == ".docx":
            doc = Document(file)
            for para in doc.paragraphs:
                text += para.text + "\n"

        elif ext == ".pptx":
            ppt = Presentation(file)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        else:
            st.warning(f"⚠️ Unsupported file format: {ext}")

    return text.strip()


# --- Split text into chunks ---
def get_text_chunks(text):
    splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = splitter.split_text(text)
    return chunks


# --- Create FAISS vector store ---
@st.cache_resource
def create_vector_store(chunks):
    embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
    store = FAISS.from_texts(chunks, embedding=embeddings)
    store.save_local("faiss_index")
    return store


# --- Build QA chain ---
def build_conversational_chain():
    prompt_template = """
    Answer the question as accurately and detailed as possible using the context below.
    If the answer is not in the context, say "The answer is not available in the context."
    Keep the tone factual and concise.

    Context:
    {context}

    Question:
    {question}

    Answer:
    """

    llm = ChatGoogleGenerativeAI(model="gemini-2.5-flash", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])

    # ✅ Create LLM chain
    llm_chain = LLMChain(llm=llm, prompt=prompt)

    # ✅ StuffDocumentsChain does the "combine all docs" logic
    chain = StuffDocumentsChain(llm_chain=llm_chain, document_variable_name="context")

    return chain


# --- Handle user query ---
def handle_user_query(question):
    embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")

    try:
        db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    except Exception:
        st.error("❌ Vector store not found. Please upload and process documents first.")
        st.stop()

    # Retrieve similar docs
    docs = db.similarity_search(question)
    chain = build_conversational_chain()

    # Run the chain
    response = chain.invoke({"input_documents": docs, "question": question})

    # ✅ Handle structured output properly
    if isinstance(response, dict):
        answer = response.get("output_text", "⚠️ No response generated.")
    else:
        answer = str(response)

    # ✅ Display in Streamlit
    st.markdown("### 💬 **Response:**")
    st.write(answer)

# --- Streamlit UI ---
def main():
    st.set_page_config(page_title="RAG Document Chat", page_icon="📘")
    st.title("📚 Multi-Document RAG-based Gemini Assistant")
    st.write("Upload PDF, DOCX, or PPTX files and ask questions based on their content.")

    with st.sidebar:
        st.header("📂 Upload Documents")
        uploaded_files = st.file_uploader(
            "Choose your documents", type=["pdf", "docx", "pptx"], accept_multiple_files=True
        )

        if st.button("📖 Process Documents"):
            if uploaded_files:
                with st.spinner("Extracting and indexing documents..."):
                    text = get_text_from_files(uploaded_files)
                    if not text.strip():
                        st.error("❌ No text could be extracted from the files.")
                        return
                    chunks = get_text_chunks(text)
                    create_vector_store(chunks)
                    st.success("✅ Documents processed and stored successfully!")
            else:
                st.warning("Please upload at least one document before processing.")

    question = st.text_input("💭 Ask a question based on your uploaded documents:")
    if question:
        handle_user_query(question)


if __name__ == "__main__":
    main()
